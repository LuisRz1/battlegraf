"""AI question generation agent using LangChain/OpenAI with local fallback.

El agente puede recibir contexto desde la memoria Redis (AgentMemory) para
evitar repetir preguntas ya generadas por el colegio (memory por school_id).
"""

import asyncio
import os
import random
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pydantic import SecretStr
from pypdf import PdfReader

from src.domain.enums import Subject
from src.domain.interfaces.repositories import QuestionAgent
from src.infrastructure.config import get_settings


class MockQuestionAgent(QuestionAgent):
    """Fallback agent that generates deterministic mock questions."""

    TEMPLATES = [
        {
            "text": "Cual es el resultado de {a} + {b}?",
            "option_a": "{a_plus_b}",
            "option_b": "{a_plus_b_wrong}",
            "option_c": "{a}",
            "option_d": "{b}",
            "correct": "A",
            "explanation": "La suma de {a} y {b} es {a_plus_b}.",
        },
        {
            "text": "Si tenemos {a} manzanas y comemos {b}, cuantas quedan?",
            "option_a": "{a_minus_b}",
            "option_b": "{a_plus_b}",
            "option_c": "{b}",
            "option_d": "{a}",
            "correct": "A",
            "explanation": "Restamos {b} de {a}, obteniendo {a_minus_b}.",
        },
        {
            "text": "Cual de estas palabras es un sinonimo de 'rapido'?",
            "option_a": "Veloz",
            "option_b": "Lento",
            "option_c": "Pesado",
            "option_d": "Tranquilo",
            "correct": "A",
            "explanation": "Veloz significa rapido.",
        },
        {
            "text": "Cual es la capital de Francia?",
            "option_a": "Paris",
            "option_b": "Londres",
            "option_c": "Madrid",
            "option_d": "Roma",
            "correct": "A",
            "explanation": "Paris es la capital de Francia.",
        },
    ]

    def _format(self, template: dict[str, str], idx: int) -> dict[str, str]:
        a = random.randint(2, 20)
        b = random.randint(1, a - 1)
        data = {
            "a": a,
            "b": b,
            "a_plus_b": a + b,
            "a_plus_b_wrong": a + b + random.randint(1, 5),
            "a_minus_b": a - b,
            "idx": idx,
        }
        return {k: v.format(**data) for k, v in template.items()}

    async def generate_questions(
        self,
        material_text: str,
        subject: Subject,
        count: int = 100,
        context: dict | None = None,
    ) -> list[dict[str, Any]]:
        """Generate deterministic mock questions."""
        _ = material_text
        _ = context
        questions = []
        for i in range(count):
            template = self.TEMPLATES[i % len(self.TEMPLATES)]
            formatted = self._format(template, i)
            questions.append(
                {
                    "subject": subject.value,
                    "text": formatted["text"],
                    "option_a": formatted["option_a"],
                    "option_b": formatted["option_b"],
                    "option_c": formatted["option_c"],
                    "option_d": formatted["option_d"],
                    "correct_option": formatted["correct"],
                    "explanation": formatted["explanation"],
                }
            )
        return questions

    async def extract_text_from_file(self, file_path: str) -> str:
        return await extract_text_from_file(file_path)


async def extract_text_from_file(file_path: str) -> str:
    """Extract text from PDF, DOCX, or TXT files using local libraries."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        doc = Document(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    if suffix == ".pptx":
        presentation = Presentation(str(path))
        return "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
    if suffix == ".txt":
        return path.read_text(encoding="utf-8")
    raise ValueError(f"Unsupported file extension: {suffix}")


class OpenAIQuestionAgent(QuestionAgent):
    """LangChain/OpenAI question generation agent."""

    def __init__(self, api_key: str | None = None, model: str | None = None,
                 base_url: str | None = None) -> None:
        settings = get_settings()
        self.api_key = (
            api_key or settings.openai_api_key or os.environ.get("OPENAI_API_KEY")
        )
        self.model = model or settings.openai_model
        self.base_url = base_url or settings.openai_base_url

    async def extract_text_from_file(self, file_path: str) -> str:
        return await extract_text_from_file(file_path)

    async def generate_questions(
        self,
        material_text: str,
        subject: Subject,
        count: int = 100,
        context: dict | None = None,
    ) -> list[dict[str, Any]]:
        """Generate questions via OpenAI or fallback to mock if no key."""
        if not self.api_key:
            return await MockQuestionAgent().generate_questions(
                material_text, subject, count, context
            )

        try:
            from langchain_openai import ChatOpenAI
        except ImportError as exc:
            raise RuntimeError("langchain-openai is not installed") from exc

        llm_kwargs = {
            "api_key": SecretStr(self.api_key),
            "model": self.model,
            "temperature": 0.7,
        }
        if self.base_url:
            llm_kwargs["base_url"] = self.base_url
        llm = ChatOpenAI(**llm_kwargs)
        prompt = self._build_prompt(material_text, subject, count, context)
        # Run blocking LangChain call in a thread pool
        response = await asyncio.to_thread(llm.invoke, prompt)
        if not isinstance(response.content, str):
            raise ValueError("Agent response must be textual JSON")
        return self._parse_response(response.content)

    def _build_prompt(
        self,
        material_text: str,
        subject: Subject,
        count: int,
        context: dict | None = None,
    ) -> str:
        memo = ""
        if context:
            prior = context.get("prior_questions") or []
            if prior:
                memo = (
                    "Evita repetir o parafrasear estas preguntas ya generadas "
                    "anteriormente para esta materia:\n"
                    + "\n".join(f"- {p[:120]}" for p in prior[:20])
                    + "\n\n"
                )
        return (
            f"Genera {count} preguntas de alternativa multiple basadas en el siguiente texto. "
            "Para cada pregunta devuelve un objeto JSON con: "
            "text, option_a, option_b, option_c, option_d, correct_option (A/B/C/D), explanation. "
            f"Materia: {subject.label}. "
            "Responde unicamente con un array JSON. No incluyas explicaciones adicionales.\n\n"
            f"{memo}"
            f"TEXTO:\n{material_text[:4000]}"
        )

    def _parse_response(self, content: str) -> list[dict[str, Any]]:
        import json

        text = content.strip()
        if text.startswith("```"):
            text = text.strip("`")
            text = text.removeprefix("json")
        try:
            return json.loads(text)
        except json.JSONDecodeError as exc:
            raise ValueError("Failed to parse agent response as JSON") from exc


def build_question_agent() -> QuestionAgent:
    """Factory that returns OpenAI agent if key is available, otherwise mock."""
    settings = get_settings()
    if settings.openai_api_key:
        return OpenAIQuestionAgent()
    return MockQuestionAgent()